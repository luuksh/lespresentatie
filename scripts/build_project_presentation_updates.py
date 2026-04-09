#!/usr/bin/env python3
from __future__ import annotations

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path("/Users/luukhijne/Desktop/Klassenplattegrond")
DATA_PATH = ROOT / "data/kerndoelen/kerndoelen-map.json"
OUT_DIR = ROOT / "updated_presentaties"

PROJECT_SOURCES = {
    "renaissance": Path("/Users/luukhijne/Library/CloudStorage/OneDrive-WillibrordStichting/CGU-AFD-Nederlands - General/4 Nederlands/Ik sta op de schouders van reuzen/Renaissance/2526/Renaissance 2526.pptx"),
    "technologie": Path("/Users/luukhijne/Library/CloudStorage/OneDrive-WillibrordStichting/CGU-AFD-Nederlands - General/4 5 Nederlands/Ik in de wereld/Vaardigheden/Lezen in thema's/Hoe je te verhouden tot technologie/Ik en technologie.pptx"),
    "netschrift": Path("/Users/luukhijne/Library/CloudStorage/OneDrive-WillibrordStichting/CGU-AFD-Nederlands - General/1 Nederlands/Netschrift/Netschrift.pptx"),
}

PROJECT_DIDACTIC_SLIDES = {
    "renaissance": [
        (
            "Renaissance | Wat Lever Je Op?",
            [
                "Je kiest met je duo een renaissanceschrijver en een primaire bron.",
                "Je onderzoekt schrijver, genre, context en waarom de tekst past bij de renaissance.",
                "Je maakt een aemulatio: een hedendaagse bewerking met behoud van functie en kern.",
                "Je presenteert auteur, bron, context en aemulatio helder en publieksgericht.",
            ],
        ),
        (
            "Renaissance | Werkroute",
            [
                "1. Kies auteur en bron.",
                "2. Begrijp elk woord en de historische context van de bron.",
                "3. Bepaal genre, functie en renaissancekenmerken.",
                "4. Maak een eigen aemulatio die inhoudelijk en vormelijk doordacht is.",
                "5. Bouw een presentatie die de klas echt helpt begrijpen wat jullie hebben gedaan.",
            ],
        ),
        (
            "Renaissance | Succescriteria",
            [
                "Leg helder uit wie de schrijver is en wat hem of haar bijzonder maakt.",
                "Zet de primaire bron zichtbaar op de PowerPoint en licht die functioneel toe.",
                "Laat zien hoe jullie aemulatio voortbouwt op het origineel.",
                "Presenteer duidelijk, verstaanbaar, publieksgericht en met een verzorgde PowerPoint.",
            ],
        ),
    ],
    "technologie": [
        (
            "Technologie | Hoofdvraag",
            [
                "Wat is volgens jou de ideale omgang met de techniek om ons heen?",
                "Je bouwt je antwoord stap voor stap op door artikelen te lezen en opdrachten te maken.",
                "Na elk artikel of gesprek stel je je standpunt bij of scherpt het aan.",
                "Je verzamelt je aantekeningen in kladschrift en werkt kernstukken uit in je netschrift.",
            ],
        ),
        (
            "Technologie | Leerroute",
            [
                "Les 1: eigen omgang met technologie en persoonlijk voornemen.",
                "Les 2: sociaal-mediagebruik, artikel lezen en standpunten wegen.",
                "Vervolg: extra artikelen vergelijken, beleid bespreken en advies formuleren.",
                "Afsluiting: schrijven over een wereld na big tech en je eigen omgang opnieuw doordenken.",
            ],
        ),
        (
            "Technologie | Waar Word Je Op Beoordeeld?",
            [
                "Je leest artikelen niet oppervlakkig, maar analyseert toon, perspectief en argumentatie.",
                "Je vergelijkt visies en weegt bruikbaarheid en overtuigingskracht af.",
                "Je verbindt die analyse aan je eigen gedrag, keuzes en standpunt.",
                "Je werkt je inzichten zorgvuldig uit in gesprek, aantekeningen en schriftelijke producten.",
            ],
        ),
    ],
    "netschrift": [
        (
            "Netschrift | Wat Is Het Doel?",
            [
                "Je netschrift is geen kladblok maar een persoonlijk groeidossier.",
                "Je laat erin zien wat je leerde, wat je herschreef en hoe je feedback gebruikte.",
                "Het schrift maakt ontwikkeling zichtbaar: inhoudelijk, talig en persoonlijk.",
                "Netheid en zorg zijn belangrijk, maar altijd in dienst van betekenis en groei.",
            ],
        ),
        (
            "Netschrift | Werkwijze",
            [
                "1. Bekijk feedback en kies sterke punten.",
                "2. Benoem concrete verbeterpunten.",
                "3. Werk feedback én feedforward uit.",
                "4. Zet een verzorgde, doordachte versie in je netschrift.",
                "5. Kijk terug: wat zegt dit over jouw ontwikkeling als schrijver en leerling?",
            ],
        ),
        (
            "Netschrift | Succescriteria",
            [
                "Je verwerkt feedback zichtbaar en eerlijk.",
                "Je formuleert nauwkeurig en verzorgt spelling en vorm.",
                "Je vat samen, ordent en kiest bewust wat belangrijk genoeg is voor je netschrift.",
                "Je maakt het schrift persoonlijk, overzichtelijk en betekenisvol voor later.",
            ],
        ),
    ],
}


def load_doc() -> dict:
    return json.loads(DATA_PATH.read_text())


def project_snapshot(doc: dict, project_id: str) -> tuple[dict, list[dict], list[dict]]:
    project = next(item for item in doc["projects"] if item["id"] == project_id)
    focus = [row for row in doc["records"] if row["projects"].get(project_id) == "focus"]
    support = [row for row in doc["records"] if row["projects"].get(project_id) == "support"]
    return project, focus, support


def add_title_and_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)
    else:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(4.8))
        frame = box.text_frame
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)


def build_bullets(project: dict, focus: list[dict], support: list[dict]) -> tuple[list[str], list[str], list[str]]:
    secondary = ", ".join(project.get("secondaryMagisterSkills", [])) or "geen"
    kompas = [
        f"Hoofdvaardigheid in Magister: {project['primaryMagisterSkill']}",
        f"Ondersteunende vaardigheden: {secondary}",
        project["assessmentSummary"],
        project["studentFacingDescription"],
    ]
    focus_bullets = [f"{row['magisterSkill']}: {row['label']}" for row in focus[:10]]
    if len(focus) > 10:
        focus_bullets.append(f"+ {len(focus) - 10} extra focuslabels in de kerndoelenkaart")
    support_bullets = [f"{row['magisterSkill']}: {row['label']}" for row in support[:8]]
    if len(support) > 8:
        support_bullets.append(f"+ {len(support) - 8} extra ondersteunende labels in de kerndoelenkaart")
    return kompas, focus_bullets, support_bullets


def update_presentation(project_id: str, source: Path, doc: dict) -> Path:
    project, focus, support = project_snapshot(doc, project_id)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    target = OUT_DIR / f"{source.stem} - bijgewerkt.pptx"
    shutil.copy2(source, target)
    prs = Presentation(target)
    kompas, focus_bullets, support_bullets = build_bullets(project, focus, support)
    for title, bullets in PROJECT_DIDACTIC_SLIDES.get(project_id, []):
        add_title_and_bullets_slide(prs, title, bullets)
    add_title_and_bullets_slide(prs, f"{project['name']} | Projectkompas", kompas)
    add_title_and_bullets_slide(prs, f"{project['name']} | Focus van de beoordeling", focus_bullets)
    add_title_and_bullets_slide(prs, f"{project['name']} | Ondersteunende labels", support_bullets)
    prs.save(target)
    return target


def main() -> int:
    doc = load_doc()
    generated = []
    for project_id, source in PROJECT_SOURCES.items():
        generated.append(update_presentation(project_id, source, doc))
    for path in generated:
        print(path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
